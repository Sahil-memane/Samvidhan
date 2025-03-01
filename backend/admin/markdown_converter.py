import os
import io
import fitz  # PyMuPDF for PDFs
import pandas as pd
from PIL import Image
# from tkinter import Tk
# from tkinter.filedialog import askopenfilename
import re
import string
import spacy
from nltk.corpus import wordnet
from docx import Document
from pptx import Presentation
import olefile

import tempfile

# Load spaCy model
nlp = spacy.load("en_core_web_sm")

# def create_folder(file_path):
#     folder_name = os.path.splitext(os.path.basename(file_path))[0]
#     if not os.path.exists(folder_name):
#         os.makedirs("output")
#     return "output"

def is_roman_numeral(text):
    roman_numeral_pattern = r'^(M{0,4}(CM|CD|D?C{0,3})(XC|XL|L?X{0,3})(IX|IV|V?I{0,3}))$'
    return re.match(roman_numeral_pattern, text.strip(), re.IGNORECASE) is not None

def has_meaning(word):
    return bool(wordnet.synsets(word))

def is_heading(text, font_size, is_italic):
    doc = nlp(text)
    heading_keywords = {"introduction", "conclusion", "summary", "abstract", "chapter", "table of contents"}
    clean_text = text.translate(str.maketrans('', '', string.punctuation)).lower()
    is_heading_keyword = any(keyword in clean_text for keyword in heading_keywords)
    
    words = clean_text.split()
    meaningful_words = [word for word in words if has_meaning(word)]
    
    is_meaningful = len(meaningful_words) / len(words) > 0.5 if words else False
    
    return (len(text) < 100 and (text.istitle() or text.isupper() or font_size > 12) and not is_roman_numeral(text) and is_meaningful) or is_heading_keyword

def preprocess_text(text_blocks):
    cleaned_lines = []
    previous_heading = None
    
    for block in text_blocks:
        text = block["text"]
        font_size = block["size"]
        is_italic = block.get("italic", False)
        text = re.sub(r'\s+', ' ', text).strip()
        if not text:
            continue
        if is_heading(text, font_size, is_italic):
            if text != previous_heading:
                cleaned_lines.append(f"\n# {text}\n")
                previous_heading = text
        else:
            cleaned_lines.append(text)
            previous_heading = None
            
    return ' '.join(cleaned_lines)

def extract_text_from_pdf(pdf_path, folder_name):
    pdf_document = fitz.open(pdf_path)
    text_blocks = []
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        blocks = page.get_text("dict")["blocks"]
        for b in blocks:
            if "lines" in b:
                for l in b["lines"]:
                    for s in l["spans"]:
                        text_blocks.append({
                            "text": s["text"],
                            "size": s["size"],
                            "italic": "italic" in s["font"]
                        })

    cleaned_text = preprocess_text(text_blocks)
    
    text_file_path = os.path.join(folder_name, "text.txt")
    markdown_file_path = os.path.join(folder_name, "text.md")
    
    with open(text_file_path, "w", encoding="utf-8") as text_file:
        text_file.write(cleaned_text)
    
    with open(markdown_file_path, "w", encoding="utf-8") as markdown_file:
        markdown_file.write(cleaned_text)

def extract_text_from_docx(docx_path, folder_name):
    doc = Document(docx_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    cleaned_text = preprocess_text([{"text": p, "size": 12, "italic": False} for p in paragraphs])
    
    text_file_path = os.path.join(folder_name, "text.txt")
    markdown_file_path = os.path.join(folder_name, "text.md")
    
    with open(text_file_path, "w", encoding="utf-8") as text_file:
        text_file.write(cleaned_text)
    
    with open(markdown_file_path, "w", encoding="utf-8") as markdown_file:
        markdown_file.write(cleaned_text)

def extract_text_from_pptx(pptx_path, folder_name):
    prs = Presentation(pptx_path)
    paragraphs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                paragraphs.append(shape.text)
    cleaned_text = preprocess_text([{"text": p, "size": 12, "italic": False} for p in paragraphs])
    
    text_file_path = os.path.join(folder_name, "text.txt")
    markdown_file_path = os.path.join(folder_name, "text.md")
    
    with open(text_file_path, "w", encoding="utf-8") as text_file:
        text_file.write(cleaned_text)
    
    with open(markdown_file_path, "w", encoding="utf-8") as markdown_file:
        markdown_file.write(cleaned_text)

def extract_text_from_ppt(ppt_path, folder_name):
    # Extract text from .ppt files using olefile
    if not olefile.isOleFile(ppt_path):
        print(f"{ppt_path} is not a valid OLE file.")
        return
    
    # Open the .ppt file
    ole = olefile.OleFileIO(ppt_path)
    text_blocks = []
    for stream in ole.listdir():
        if stream[0] == 'PowerPoint Document':
            data = ole.openstream(stream)
            text = data.read().decode('utf-8', errors='ignore')
            text_blocks.append({"text": text, "size": 12, "italic": False})
    
    cleaned_text = preprocess_text(text_blocks)
    
    text_file_path = os.path.join(folder_name, "text.txt")
    markdown_file_path = os.path.join(folder_name, "text.md")
    
    with open(text_file_path, "w", encoding="utf-8") as text_file:
        text_file.write(cleaned_text)
    
    with open(markdown_file_path, "w", encoding="utf-8") as markdown_file:
        markdown_file.write(cleaned_text)

def extract_images_from_pdf(pdf_path, folder_name):
    pdf_document = fitz.open(pdf_path)
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image = Image.open(io.BytesIO(image_bytes))
            image.save(os.path.join(folder_name, f"image_page_{page_num+1}_{img_index+1}.{image_ext}"))

# def main():
#     Tk().withdraw()
#     file_path = askopenfilename(filetypes=[("All files", "*.*")])
#     if not file_path:
#         print("No file selected.")
#         return
    
#     file_ext = os.path.splitext(file_path)[1].lower()
#     folder_name = create_folder(file_path)
    
#     if file_ext == ".pdf":
#         extract_text_from_pdf(file_path, folder_name)
#         extract_images_from_pdf(file_path, folder_name)
#     elif file_ext == ".docx":
#         extract_text_from_docx(file_path, folder_name)
#     elif file_ext == ".pptx":
#         extract_text_from_pptx(file_path, folder_name)
#     elif file_ext == ".ppt":
#         extract_text_from_ppt(file_path, folder_name)
#     else:
#         print("Unsupported file format.")
#         return
    
#     print(f"Extraction completed. Check the folder: {folder_name}")

# if __name__ == "__main__":
#     main()

def convert_to_markdown(file_path):
    print(file_path)
    
    file_ext = os.path.splitext(file_path)[1].lower()
    folder_name = "output"
    
    if file_ext == ".pdf":
        extract_text_from_pdf(file_path, folder_name)
        extract_images_from_pdf(file_path, folder_name)
    elif file_ext == ".docx":
        extract_text_from_docx(file_path, folder_name)
    elif file_ext == ".pptx":
        extract_text_from_pptx(file_path, folder_name)
    elif file_ext == ".ppt":
        extract_text_from_ppt(file_path, folder_name)
    else:
        return False
    
    return True
    
    
    
    
